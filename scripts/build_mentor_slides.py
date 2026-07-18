"""Generate a 16:9 PPTX deck for the FinAdvisor mentor sessions.

Output: financial-advisor-llm/MentorSessions.pptx

Upload to Google Drive and "Open with Google Slides" to convert natively.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "MentorSessions.pptx"

NAVY = RGBColor(0x0D, 0x2B, 0x4E)
TEAL = RGBColor(0x14, 0x8F, 0x9B)
ORANGE = RGBColor(0xE7, 0x7E, 0x22)
GREY = RGBColor(0x55, 0x5C, 0x66)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is None:
        shape.line.fill.background()
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=BLACK, align="left", font="Calibri"):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=BLACK,
                bullet_color=None, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    bc = bullet_color or TEAL
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bc
        bullet.font.name = font
        body = p.add_run()
        body.text = item
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = font
    return tb


def header_bar(slide, session_label, session_color):
    add_rect(slide, 0, 0, SW, Inches(0.55), NAVY)
    add_text(slide, Inches(0.4), Inches(0.10), Inches(8), Inches(0.4),
             "FinAdvisor — Mentor Sessions", size=16, bold=True, color=WHITE)
    chip_w = Inches(2.4)
    chip = add_rect(slide, SW - chip_w - Inches(0.4), Inches(0.10),
                    chip_w, Inches(0.35), session_color)
    add_text(slide, SW - chip_w - Inches(0.4), Inches(0.13), chip_w, Inches(0.3),
             session_label, size=12, bold=True, color=WHITE, align="center")


def footer_bar(slide, text):
    add_rect(slide, 0, SH - Inches(0.35), SW, Inches(0.35), LIGHT)
    add_text(slide, Inches(0.4), SH - Inches(0.3), SW - Inches(0.8), Inches(0.3),
             text, size=11, color=GREY, align="left")


def title_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(2.4), SW, Inches(0.05), TEAL)
    add_text(s, Inches(0.8), Inches(1.0), Inches(12), Inches(0.6),
             "FinAdvisor", size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(1.65), Inches(12), Inches(0.5),
             "Personalized Financial Advisor using LLM",
             size=24, color=RGBColor(0xCF, 0xE3, 0xE6))
    add_text(s, Inches(0.8), Inches(2.7), Inches(12), Inches(0.5),
             "4-Session Mentor Plan",
             size=28, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.3), Inches(12), Inches(0.5),
             "Capstone — Advanced Certification in Agentic & Generative AI",
             size=16, color=RGBColor(0xCF, 0xE3, 0xE6))
    add_text(s, Inches(0.8), Inches(5.6), Inches(12), Inches(0.4),
             "Reference paper: FinGPT (Yang, Liu, Wang 2023)",
             size=14, color=RGBColor(0xCF, 0xE3, 0xE6))
    add_text(s, Inches(0.8), Inches(6.0), Inches(12), Inches(0.4),
             "Stack: Llama-3.3-70B (HF→Groq) · Alpha Vantage · Chroma RAG · Streamlit",
             size=14, color=RGBColor(0xCF, 0xE3, 0xE6))
    add_text(s, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
             "Timeline: 2026-07-04 → 2026-07-25 · 4 weekly mentored sessions",
             size=14, color=ORANGE, bold=True)


def overview_slide():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "OVERVIEW", TEAL)
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "The 4-session arc — progressive reveal", size=26, bold=True,
             color=NAVY)
    add_text(s, Inches(0.4), Inches(1.5), Inches(12), Inches(0.4),
             "Each session adds one major capability. Same prompts get smarter as layers stack.",
             size=14, color=GREY)
    sessions = [
        ("Session 1", "Problem framed,\nbaseline talking",
         "Open-weights LLM call\n+ live AV data\n(no grounding yet)",
         "Notebook:\nAAPL → LLM paragraph",
         RGBColor(0x4A, 0x90, 0xE2)),
        ("Session 2", "Grounded answers\nvia RAG",
         "Hybrid retrieval\n(BM25 + dense + RRF)\nwith source citations",
         "Notebook:\nRoth IRA query w/ sources",
         RGBColor(0x14, 0x8F, 0x9B)),
        ("Session 3", "Agent: tools +\nreasoning + persona",
         "ReAct loop, tool registry,\nSQLite profile,\nStreamlit UI",
         "3 prompts:\nconcept / live / goal",
         RGBColor(0xE7, 0x7E, 0x22)),
        ("Session 4", "Polished product\n+ eval + safety",
         "Full app, eval table,\nred-team results,\nfuture-work bridge",
         "Final defense:\ndeployed URL + numbers",
         RGBColor(0xC0, 0x39, 0x2B)),
    ]
    card_w = Inches(3.0)
    card_h = Inches(4.5)
    gap = Inches(0.13)
    total = card_w * 4 + gap * 3
    start_x = (SW - total) // 2
    for i, (label, theme, whats_new, artifact, color) in enumerate(sessions):
        x = start_x + (card_w + gap) * i
        y = Inches(2.1)
        add_rect(s, x, y, card_w, card_h, LIGHT)
        add_rect(s, x, y, card_w, Inches(0.5), color)
        add_text(s, x, y + Inches(0.08), card_w, Inches(0.4),
                 label, size=14, bold=True, color=WHITE, align="center")
        add_text(s, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4),
                 Inches(1.0), theme, size=15, bold=True, color=NAVY)
        add_text(s, x + Inches(0.2), y + Inches(1.85), card_w - Inches(0.4),
                 Inches(0.3), "WHAT'S NEW", size=10, bold=True, color=GREY)
        add_text(s, x + Inches(0.2), y + Inches(2.15), card_w - Inches(0.4),
                 Inches(1.4), whats_new, size=12, color=BLACK)
        add_text(s, x + Inches(0.2), y + Inches(3.55), card_w - Inches(0.4),
                 Inches(0.3), "ANCHOR ARTIFACT", size=10, bold=True, color=GREY)
        add_text(s, x + Inches(0.2), y + Inches(3.85), card_w - Inches(0.4),
                 Inches(0.6), artifact, size=12, color=BLACK)
    footer_bar(s, "Same MSFT and retirement prompts reappear across sessions — the answer visibly gets smarter.")


def session_title_slide(num, title, subtitle, color):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, color)
    add_text(s, Inches(0.8), Inches(2.5), Inches(12), Inches(0.6),
             f"SESSION {num}", size=20, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, Inches(0.8), Inches(3.0), Inches(12), Inches(1.0),
             title, size=42, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.6), Inches(12), Inches(0.6),
             subtitle, size=20, color=RGBColor(0xEC, 0xEF, 0xF1))


def session_overview(num, theme_color, theme, demo_minutes, anchor):
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {num} · OVERVIEW", theme_color)
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Theme & demo at a glance", size=26, bold=True, color=NAVY)
    add_rect(s, Inches(0.4), Inches(1.7), Inches(6.1), Inches(2.3), LIGHT)
    add_text(s, Inches(0.6), Inches(1.85), Inches(5.8), Inches(0.4),
             "ONE-LINE THEME", size=11, bold=True, color=GREY)
    add_text(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.7),
             theme, size=16, color=BLACK)
    add_rect(s, Inches(6.7), Inches(1.7), Inches(6.2), Inches(2.3), LIGHT)
    add_text(s, Inches(6.9), Inches(1.85), Inches(5.8), Inches(0.4),
             "DEMO LENGTH", size=11, bold=True, color=GREY)
    add_text(s, Inches(6.9), Inches(2.2), Inches(5.8), Inches(0.6),
             demo_minutes, size=22, bold=True, color=theme_color)
    add_text(s, Inches(6.9), Inches(2.95), Inches(5.8), Inches(0.4),
             "ANCHOR ARTIFACT", size=11, bold=True, color=GREY)
    add_text(s, Inches(6.9), Inches(3.3), Inches(5.8), Inches(0.6),
             anchor, size=14, color=BLACK)
    return s


def two_col_slide(num, theme_color, title, left_h, left_items, right_h, right_items,
                  footer=""):
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {num}", theme_color)
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=NAVY)
    add_rect(s, Inches(0.4), Inches(1.65), Inches(6.1), Inches(0.5), theme_color)
    add_text(s, Inches(0.55), Inches(1.72), Inches(6), Inches(0.4),
             left_h, size=14, bold=True, color=WHITE)
    add_bullets(s, Inches(0.55), Inches(2.3), Inches(5.95), Inches(4.5),
                left_items, size=14, bullet_color=theme_color)
    add_rect(s, Inches(6.7), Inches(1.65), Inches(6.2), Inches(0.5), NAVY)
    add_text(s, Inches(6.85), Inches(1.72), Inches(6), Inches(0.4),
             right_h, size=14, bold=True, color=WHITE)
    add_bullets(s, Inches(6.85), Inches(2.3), Inches(6.05), Inches(4.5),
                right_items, size=14, bullet_color=NAVY)
    if footer:
        footer_bar(s, footer)


def list_slide(num, theme_color, title, header, items, footer=""):
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {num}", theme_color)
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.4), Inches(1.5), Inches(12), Inches(0.4),
             header, size=14, color=GREY)
    add_bullets(s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(5.0),
                items, size=15, bullet_color=theme_color)
    if footer:
        footer_bar(s, footer)


SESSIONS_META = [
    {
        "num": 1, "color": RGBColor(0x4A, 0x90, 0xE2),
        "title": "Problem framed, baseline talking",
        "subtitle": "Open-weights LLM + live data — nothing grounded yet",
        "theme": ("We can talk to an open-weights LLM and pull live finance "
                  "data. Nothing is grounded yet — that's the point."),
        "demo_min": "5 minutes",
        "anchor": "Live notebook: AAPL quote → LLM paragraph",
    },
    {
        "num": 2, "color": RGBColor(0x14, 0x8F, 0x9B),
        "title": "Grounded answers via RAG",
        "subtitle": "Every claim citable to a source",
        "theme": ("Now every answer can cite a source. The corpus is real, "
                  "the retrieval is hybrid, the grounding is visible."),
        "demo_min": "8 minutes",
        "anchor": "Notebook: \"Roth IRA?\" with [Source: …] citations",
    },
    {
        "num": 3, "color": RGBColor(0xE7, 0x7E, 0x22),
        "title": "Agent: tools + reasoning + personalization",
        "subtitle": "The LLM picks its own tools and uses your profile",
        "theme": ("Now the LLM decides for itself when to call AV, "
                  "when to call a calculator, and uses your profile."),
        "demo_min": "10 minutes",
        "anchor": "3 scripted prompts: concept / live / goal",
    },
    {
        "num": 4, "color": RGBColor(0xC0, 0x39, 0x2B),
        "title": "Polished product + eval + safety",
        "subtitle": "Working product, measured quality, honest limits",
        "theme": ("It's a working product, with numbers proving quality, "
                  "and we're honest about what's left."),
        "demo_min": "12 minutes",
        "anchor": "Full Streamlit demo + eval table + red-team results",
    },
]


def build_session_1():
    m = SESSIONS_META[0]
    session_title_slide(m["num"], m["title"], m["subtitle"], m["color"])
    session_overview(m["num"], m["color"], m["theme"], m["demo_min"], m["anchor"])

    two_col_slide(
        m["num"], m["color"], "Demo script (5 min)",
        "DEMO STEPS",
        [
            "Open the proposal in Google Docs — walk sections 1-4 in 90s",
            "get_quote(\"AAPL\") → JSON; point at cache hit",
            "Paste quote into chat() with system prompt → coherent paragraph",
            "Repeat with get_company_overview(\"AAPL\") for richer answer",
            "Pose the gap: \"Where did the P/E claim come from?\"",
        ],
        "WHAT TO HIGHLIGHT",
        [
            "One credential (HF_TOKEN), provider routes to Groq",
            "Cache hit → instant; AV free tier = 25 req/day",
            "System prompt enforces educational framing + disclaimer",
            "No memory yet, no grounding yet, no tools yet",
            "This is the foundation — Session 2 adds grounding",
        ],
        footer="Files: src/advisor/llm/client.py · src/advisor/tools/alpha_vantage.py · src/advisor/llm/prompts.py",
    )

    talking = [
        ("Why open-weights via Hugging Face Inference Providers",
         "FinGPT thesis is open-source financial LLMs. HF Inference "
         "Providers gives one credential (HF_TOKEN), free tier, "
         "OpenAI-compatible chat. HF transparently routes to Groq for "
         "speed; switching providers is one .env line. The code never "
         "imports a provider-specific SDK — architecture is reversible."),
        ("Why Alpha Vantage (not yfinance / Polygon / IEX)",
         "AV is what the spec calls for, and it's the right cost/benefit "
         "point: quotes + overviews + news/sentiment + technicals + "
         "sectors + FX from one API. 25 req/day is tight, so we use "
         "requests-cache with a 30-min TTL. Cache pre-warms the demo so "
         "we never blow quota mid-presentation."),
        ("Why we deferred fine-tuning, and why that's still aligned with FinGPT",
         "FinGPT shows LoRA makes domain adaptation cheap (~$300). We "
         "don't have GPU budget in 4 weeks, so we take the inference-side "
         "path: prompt + RAG + tools. The two paths compose — Session 4 "
         "future-work is exactly LoRA-tuning the base model and re-eval "
         "against the same benchmarks, leaving everything else identical."),
        ("Why a single Settings object reads from .env",
         "Every module imports `settings`, never os.environ. One source "
         "of truth, type-checked by pydantic-settings. .env is gitignored. "
         "Credential hygiene is a baseline safety property — relevant in "
         "Session 4 when discussing production safety."),
    ]
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {m['num']} · TALKING POINTS", m["color"])
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Talking points (elaborated)", size=26, bold=True, color=NAVY)
    y = Inches(1.7)
    for label, body in talking:
        add_text(s, Inches(0.4), y, Inches(12.5), Inches(0.4),
                 label, size=14, bold=True, color=m["color"])
        add_text(s, Inches(0.4), y + Inches(0.4), Inches(12.5), Inches(1.0),
                 body, size=11, color=BLACK)
        y += Inches(1.35)

    list_slide(
        m["num"], m["color"],
        "Limitations & expected mentor questions",
        "Volunteer the gaps before they're asked — sets up Session 2.",
        [
            "LIMITATION  ·  No memory between turns yet.",
            "LIMITATION  ·  Nothing grounded — could fabricate a P/E ratio.",
            "LIMITATION  ·  AV's 25 req/day ceiling is real (cached around it).",
            "Q  ·  \"Why not just use ChatGPT's API?\"  →  Open-weights thesis; provider-agnostic abstraction is reversible.",
            "Q  ·  \"What sources should ground the answers?\"  →  That's exactly Session 2.",
            "Q  ·  \"How do you handle rate limits?\"  →  30-min cache TTL; cache pre-warmed before demos.",
        ],
        footer="Pre-session: notebook executed top-to-bottom · cache pre-warmed for AAPL/MSFT/VOO · placeholder .env",
    )


def build_session_2():
    m = SESSIONS_META[1]
    session_title_slide(m["num"], m["title"], m["subtitle"], m["color"])
    session_overview(m["num"], m["color"], m["theme"], m["demo_min"], m["anchor"])

    two_col_slide(
        m["num"], m["color"], "Demo script (8 min)",
        "DEMO STEPS",
        [
            "Run `make index` — show Chroma directory + chunk count",
            "search(\"What is a Roth IRA?\", k=4) — top-4 with sources, BM25 vs dense scores",
            "BM25 wins: \"VOO expense ratio\" (exact ticker)",
            "Dense wins: \"safe investments for retirees\" (no token overlap)",
            "Apple question with vs without RAG — same prompt, demonstrably more grounded",
        ],
        "WHAT TO HIGHLIGHT",
        [
            "Inline [Source: glossary.md] citations in the LLM output",
            "RRF (k=60) merges two ranked lists, no score normalization needed",
            "Local Chroma → single SQLite file, survives restarts, gitignored",
            "Public corpus (SEC + IRS + FINRA + glossary) — auditable, free",
            "Source attribution > raw retrieval accuracy for trust",
        ],
        footer="Files: src/advisor/rag/ingest.py · src/advisor/rag/retrieve.py · src/advisor/rag/store.py",
    )

    talking = [
        ("Why hybrid retrieval beats dense-only",
         "Dense (BGE-small) handles semantic similarity — \"retirement "
         "income\" matches IRA passages without literal overlap. But "
         "finance is acronyms (MSFT, ETF, RMD, AGI) where lexical match "
         "matters more. BM25 nails those. RRF (Cormack 2009, k=60) merges "
         "ranks without score normalization. Robust on both query shapes."),
        ("Why local Chroma, not Pinecone or pgvector",
         "PersistentClient writes chunks + embeddings + metadata to a "
         "single SQLite file. No service, no infra, no API key, survives "
         "restarts, gitignored. For a capstone with hundreds of chunks "
         "this is correct. Production-at-scale would warrant Pinecone, "
         "but the abstraction is generic — swap is contained."),
        ("Corpus design — public, auditable, free",
         "SEC investor.gov, IRS publications (590 / 17 / 525), FINRA "
         "investor education, custom glossary. All public domain or "
         "open-license, all auditable (mentor can open the source PDF), "
         "all free. Seed index has 8 chunks today; production target is "
         "200+. Char-based chunking (size 800, overlap 120) is good "
         "enough for English finance text."),
        ("Why source attribution matters more than raw retrieval accuracy",
         "Finance failure mode is \"wrong + we couldn't tell.\" Inline "
         "[Source: …] makes every claim auditable. If the LLM cites, the "
         "user can verify. If it doesn't, that's a signal to be skeptical. "
         "Session 4's safety layer double-checks that advisory responses "
         "include source/tool markers."),
    ]
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {m['num']} · TALKING POINTS", m["color"])
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Talking points (elaborated)", size=26, bold=True, color=NAVY)
    y = Inches(1.7)
    for label, body in talking:
        add_text(s, Inches(0.4), y, Inches(12.5), Inches(0.4),
                 label, size=14, bold=True, color=m["color"])
        add_text(s, Inches(0.4), y + Inches(0.4), Inches(12.5), Inches(1.0),
                 body, size=11, color=BLACK)
        y += Inches(1.35)

    list_slide(
        m["num"], m["color"],
        "Limitations & expected mentor questions",
        "Honest about scope — sets up Session 3 (the LLM acts, not just reads).",
        [
            "LIMITATION  ·  Corpus is small (~8 chunks today; production needs hundreds).",
            "LIMITATION  ·  Char-based chunking; token-aware is a future improvement.",
            "LIMITATION  ·  No reranker yet (BGE-reranker-base would tighten top-k).",
            "LIMITATION  ·  LLM still doesn't act — it can only read what RAG hands it.",
            "Q  ·  \"Is the corpus comprehensive?\"  →  Not yet; sized for demo. Production needs continuous ingestion.",
            "Q  ·  \"What if the source is wrong?\"  →  SEC/IRS/FINRA are authoritative for US retail finance.",
            "Q  ·  \"How does the LLM use live data?\"  →  That's exactly Session 3.",
        ],
        footer="Pre-session: corpus expanded · `make index` clean from fresh dir · BM25 vs dense win examples saved",
    )


def build_session_3():
    m = SESSIONS_META[2]
    session_title_slide(m["num"], m["title"], m["subtitle"], m["color"])
    session_overview(m["num"], m["color"], m["theme"], m["demo_min"], m["anchor"])

    two_col_slide(
        m["num"], m["color"], "Demo script — 3 scripted prompts (10 min)",
        "PROMPT TYPES",
        [
            "CONCEPT  ·  \"Explain dollar-cost averaging.\"  →  RAG only, 0 tool calls",
            "LIVE     ·  \"Should I be worried about MSFT?\"  →  parallel get_company_overview + get_news_sentiment",
            "GOAL     ·  \"30, want $2M at 65, save $1500/mo — on track?\"  →  retirement_projection",
            "Surface the trace each time: which tools fired, what came back",
            "ReAct loop is ~80 lines — walk through it line by line",
        ],
        "WHAT TO HIGHLIGHT",
        [
            "MAX_STEPS = 6 (enough headroom; prevents runaway loops)",
            "Tool output truncated to 4000 chars (context safety)",
            "OpenAI-compatible tool schemas → provider-agnostic",
            "SQLite profile → \"USER CONTEXT\" block in system prompt",
            "Calculators echo their assumptions (assumed_return, years)",
        ],
        footer="Files: src/advisor/agent/react.py · tools/registry.py · tools/calculators.py · agent/memory.py",
    )

    talking = [
        ("Why ReAct (Reason + Act)",
         "ReAct (Yao et al. 2023) is the simplest agent pattern that "
         "actually works for tool use. Model reasons, optionally requests "
         "tools, agent runs them, results feed back, loop until a final "
         "answer with no tool calls. Transparent — every step is visible "
         "in the message list — exactly what we want for a capstone "
         "defense. Frameworks (LangGraph, smolagents) wrap this same "
         "pattern with extra abstractions."),
        ("Why a custom loop, not LangGraph or smolagents",
         "LangGraph is right for production agents with branching, "
         "persistence, complex state machines — overkill for a single-turn "
         "ReAct. We picked custom for transparency. \"What happens when a "
         "tool throws?\" → I point at try/except in react.py. With a "
         "framework, the answer is \"the framework handles it\" — fine "
         "for production, weak for a capstone."),
        ("Why deterministic calculators instead of LLM arithmetic",
         "Two failure modes. (1) Hallucination: 70B models drift 30%+ on "
         "long compounding formulas. (2) Wrong assumptions: 8% retail-blog "
         "default vs 7% conservative. Calculators always return the right "
         "number AND echo their assumptions back so user + LLM both see "
         "them. LLM picks the *what*; the function computes the *how*."),
        ("SQLite profile, not Redis or in-memory",
         "Loaded once per agent call into the system prompt as a USER "
         "CONTEXT block. User never re-types \"30, moderate risk\". One "
         "file, no service, survives restarts, easy to inspect. "
         "Multi-tenant production would use Postgres — abstraction in "
         "memory.py is generic, swap doesn't touch other modules."),
    ]
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {m['num']} · TALKING POINTS", m["color"])
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Talking points (elaborated)", size=26, bold=True, color=NAVY)
    y = Inches(1.7)
    for label, body in talking:
        add_text(s, Inches(0.4), y, Inches(12.5), Inches(0.4),
                 label, size=14, bold=True, color=m["color"])
        add_text(s, Inches(0.4), y + Inches(0.4), Inches(12.5), Inches(1.0),
                 body, size=11, color=BLACK)
        y += Inches(1.35)

    list_slide(
        m["num"], m["color"],
        "Limitations & expected mentor questions",
        "Limits called out — sets up Session 4 (eval numbers prove quality).",
        [
            "LIMITATION  ·  Open-weights tool-calling occasionally malforms args (caught + reported).",
            "LIMITATION  ·  No streaming yet (deferred to Session 4 stretch).",
            "LIMITATION  ·  One tenant (demo-user); multi-tenant straightforward but unbuilt.",
            "LIMITATION  ·  Step budget can be hit on pathological prompts (clean error).",
            "Q  ·  \"What if it picks the wrong tool?\"  →  It happens; eval numbers in Session 4 quantify it.",
            "Q  ·  \"Why MAX_STEPS=6?\"  →  Empirically enough; prevents runaways.",
            "Q  ·  \"How do you prove it's safe and accurate?\"  →  Exactly Session 4.",
        ],
        footer="Pre-session: 3 prompts as runnable script · pre-warmed AV cache · profile pre-populated · pytest -q green",
    )


def build_session_4():
    m = SESSIONS_META[3]
    session_title_slide(m["num"], m["title"], m["subtitle"], m["color"])
    session_overview(m["num"], m["color"], m["theme"], m["demo_min"], m["anchor"])

    two_col_slide(
        m["num"], m["color"], "Demo script (12 min)",
        "STREAMLIT WALK-THROUGH",
        [
            "Sidebar  ·  Set persona (30, moderate, $80k, retirement)",
            "Chat     ·  MSFT question (same as Session 3) — grounded + cited + disclaimed",
            "Portfolio·  Paste 3 holdings → live quotes → pie / bar → LLM commentary",
            "Goals    ·  Retirement projection with Plotly chart",
            "Markets  ·  Sector heatmap + news sentiment feed",
        ],
        "EVAL + SAFETY SLIDE",
        [
            "Financial PhraseBank accuracy on N=200 vs FinGPT-paper baselines",
            "10-task LLM-as-judge: factual / grounded / personalized / safe (1-5)",
            "Safety red-team: directive ask, guaranteed-profit, prompt injection, distress",
            "safety.py: input flagging + scrub + disclaimer enforcement",
            "Architecture sequence diagram (UI → Agent → RAG/LLM/Tools → Safety → UI)",
        ],
        footer="Files: app/streamlit_app.py · app/pages/*.py · src/advisor/eval/*.py · src/advisor/agent/safety.py",
    )

    talking = [
        ("Why Streamlit, not FastAPI + React",
         "Capstone question is \"how fast can a polished, multi-page UI "
         "ship.\" Streamlit's answer: one Python file per page, persistent "
         "state via session_state, plotly out of the box, deploys to HF "
         "Spaces in 5 min. FastAPI + React = 5-10x engineering, would NOT "
         "improve eval, safety, or agent quality. Trade-off accepted: "
         "Streamlit reruns the script on every interaction → "
         "@st.cache_resource for the retriever."),
        ("What each eval number proves — and what it doesn't",
         "FPB: model can read finance text. NOT \"gives good advice.\" "
         "10-task LLM-judge probes end-to-end quality across the 3 "
         "pillars but has the standard judge biases (favors verbose). "
         "Safety red-team: 4 attacks; pass rate measures whether "
         "disclaimer fires + directive language scrubbed + injection "
         "rebuffed. Does NOT prove robust safety under arbitrary attack "
         "— that needs a moderation classifier."),
        ("The safety architecture",
         "Three lightweight, transparent guards. (1) check_user_input — "
         "flags injection patterns + distress signals (flags, never "
         "blocks). (2) scrub_directive_language — regex over LLM output "
         "to remove \"guaranteed\" / \"risk-free\" / \"must buy\". "
         "(3) enforce_disclaimer — appends if missing. All via "
         "post_process(). Honest scope: production needs a refusal "
         "model + content moderation API; we don't pretend to have those."),
        ("The FinGPT bridge — natural future-work",
         "LoRA fine-tune the base model on Financial PhraseBank, re-run "
         "the same eval, compare. Inference layer unchanged → fine-tune "
         "is a plug-in, not a rewrite. That's the bridge: this project "
         "is the inference-side path PREPARED for the training-side "
         "extension. Other future work: moderation classifier, backtest "
         "tool, token streaming, multi-tenant memory, larger corpus."),
    ]
    s = prs.slides.add_slide(BLANK)
    header_bar(s, f"SESSION {m['num']} · TALKING POINTS", m["color"])
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Talking points (elaborated)", size=26, bold=True, color=NAVY)
    y = Inches(1.7)
    for label, body in talking:
        add_text(s, Inches(0.4), y, Inches(12.5), Inches(0.4),
                 label, size=14, bold=True, color=m["color"])
        add_text(s, Inches(0.4), y + Inches(0.4), Inches(12.5), Inches(1.0),
                 body, size=11, color=BLACK)
        y += Inches(1.35)

    list_slide(
        m["num"], m["color"],
        "Limitations & future work — closing slide",
        "Volunteer the limits up-front. Position FinGPT LoRA path as the natural extension.",
        [
            "LIMITATION  ·  Corpus is small relative to production deployment.",
            "LIMITATION  ·  AV free-tier quota constrains live data freshness.",
            "LIMITATION  ·  No fine-tune (called out as the natural extension).",
            "LIMITATION  ·  Safety = regex + disclaimer; production needs a classifier.",
            "LIMITATION  ·  Single-user; multi-tenant straightforward but unbuilt.",
            "FUTURE  ·  LoRA on Financial PhraseBank + re-eval (the FinGPT bridge).",
            "FUTURE  ·  Moderation classifier replacing regex safety.",
            "FUTURE  ·  Backtest tool + token streaming + multi-tenant + larger corpus.",
        ],
        footer="Pre-session: HF Space deployed · `make eval` results.json · screencast backup · final report 48h ahead",
    )


def pacing_slide():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "PACING TIPS", TEAL)
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Cross-session pacing tips", size=26, bold=True, color=NAVY)
    items = [
        "Sessions 1-2 are notebook-heavy — show code visibly. Mentors evaluating an engineering capstone want to see the engineering.",
        "Sessions 3-4 are product-heavy — show the Streamlit app and the eval numbers. By Session 4 the mentor wants the result, not the lines.",
        "Reuse the same scripted prompts. The MSFT question and the retirement question reappear with progressively richer answers. That delta is your most persuasive arc.",
        "End each session with a question that motivates the next. Turns a status update into a serial narrative.",
        "Always close with limitations. Volunteering them before the mentor asks signals you understand your own system.",
        "Have a backup screencast. If a live demo fails (AV down, model 5xx), switching to recorded video keeps the session on track.",
    ]
    add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.2), Inches(5.2),
                items, size=15, bullet_color=ORANGE)


def gantt_slide():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "TIMELINE", TEAL)
    add_text(s, Inches(0.4), Inches(0.85), Inches(12), Inches(0.6),
             "Mapping back to the proposal", size=26, bold=True, color=NAVY)
    add_text(s, Inches(0.4), Inches(1.5), Inches(12), Inches(0.4),
             "4 weeks · 4 mentored sessions · 2026-07-04 → 2026-07-25",
             size=14, color=GREY)
    rows = [
        ("Week", "Date", "Session", "Demo content"),
        ("1", "2026-07-10", "Session 1", "LLM client + AV wrappers + baseline notebook"),
        ("2", "2026-07-17", "Session 2", "RAG corpus + hybrid retrieval + grounded answers"),
        ("3", "2026-07-24", "Session 3", "ReAct agent + tools + Streamlit UI + 3 prompts"),
        ("4", "2026-07-25", "Session 4 — Final Defense",
         "Deployed app + eval numbers + safety + future work"),
    ]
    col_widths = [Inches(1.0), Inches(2.0), Inches(2.8), Inches(6.7)]
    x_start = Inches(0.4)
    y_start = Inches(2.2)
    row_h = Inches(0.7)
    colors = [
        WHITE,
        RGBColor(0x4A, 0x90, 0xE2),
        RGBColor(0x14, 0x8F, 0x9B),
        RGBColor(0xE7, 0x7E, 0x22),
        RGBColor(0xC0, 0x39, 0x2B),
    ]
    for ri, row in enumerate(rows):
        x = x_start
        bg = NAVY if ri == 0 else (LIGHT if ri % 2 else WHITE)
        text_color = WHITE if ri == 0 else BLACK
        for ci, cell in enumerate(row):
            add_rect(s, x, y_start + row_h * ri, col_widths[ci], row_h, bg)
            if ci == 2 and ri > 0:
                chip_w = Inches(0.3)
                chip = add_rect(s, x + Inches(0.15), y_start + row_h * ri + Inches(0.15),
                                chip_w, Inches(0.4), colors[ri])
                add_text(s, x + Inches(0.55),
                         y_start + row_h * ri + Inches(0.20),
                         col_widths[ci] - Inches(0.5), Inches(0.4),
                         cell, size=12, bold=True, color=text_color)
            else:
                add_text(s, x + Inches(0.15),
                         y_start + row_h * ri + Inches(0.20),
                         col_widths[ci] - Inches(0.2), Inches(0.4),
                         cell, size=12, bold=(ri == 0), color=text_color)
            x += col_widths[ci]
    add_text(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.5),
             "Sessions 3 and 4 are back-to-back. S3 = engineering review (does it work end-to-end?). S4 = defense (production + measured).",
             size=12, color=GREY)


def closing_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    add_rect(s, 0, Inches(2.4), SW, Inches(0.05), TEAL)
    add_text(s, Inches(0.8), Inches(2.7), Inches(12), Inches(0.8),
             "Thank you", size=48, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.6), Inches(12), Inches(0.5),
             "Questions, feedback, and direction welcome.",
             size=20, color=RGBColor(0xCF, 0xE3, 0xE6))
    add_text(s, Inches(0.8), Inches(5.8), Inches(12), Inches(0.4),
             "Repository: financial-advisor-llm/",
             size=14, color=RGBColor(0xCF, 0xE3, 0xE6))
    add_text(s, Inches(0.8), Inches(6.2), Inches(12), Inches(0.4),
             "Reference: Yang, Liu, Wang. FinGPT — arXiv:2306.06031 (2023)",
             size=14, color=RGBColor(0xCF, 0xE3, 0xE6))


title_slide()
overview_slide()
build_session_1()
build_session_2()
build_session_3()
build_session_4()
pacing_slide()
gantt_slide()
closing_slide()

prs.save(OUT)
print(f"Wrote {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")
